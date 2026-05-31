"""
T-Pot 허니팟 데이터셋 분석 PPT 생성기
python-pptx + matplotlib 사용
"""

import io
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
import numpy as np
import pandas as pd
from pptx import Presentation

# WSL2 Windows 폰트에서 한국어 폰트 등록
_KR_FONT_PATHS = [
    "/mnt/c/Windows/Fonts/malgun.ttf",
    "/mnt/c/Windows/Fonts/NotoSansKR-VF.ttf",
]
_kr_font = None
for _fp in _KR_FONT_PATHS:
    if Path(_fp).exists():
        fm.fontManager.addfont(_fp)
        _prop = fm.FontProperties(fname=_fp)
        _kr_font = _prop.get_name()
        break

if _kr_font:
    plt.rcParams["font.family"] = _kr_font
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── 색상 팔레트 ─────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # 짙은 네이비
C_ACCENT   = RGBColor(0x00, 0xC8, 0xFF)   # 사이버 블루
C_ACCENT2  = RGBColor(0xFF, 0x6B, 0x35)   # 오렌지
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD     = RGBColor(0x17, 0x2A, 0x3F)   # 카드 배경

HEX_BG    = "#0D1B2A"
HEX_ACC   = "#00C8FF"
HEX_ACC2  = "#FF6B35"
HEX_WHITE = "#FFFFFF"
HEX_LIGHT = "#B0C4DE"
HEX_CARD  = "#172A3F"
HEX_GRID  = "#1E3A50"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT_DIR = Path(__file__).parent
PPT_PATH = OUT_DIR / "tpot_dataset_analysis.pptx"
CSV_DIR  = OUT_DIR / "csv"


# ── matplotlib 공통 설정 ──────────────────────────────────────────────────────
plt.rcParams.update({
    "figure.facecolor": HEX_BG,
    "axes.facecolor":   HEX_CARD,
    "axes.edgecolor":   HEX_GRID,
    "axes.labelcolor":  HEX_WHITE,
    "xtick.color":      HEX_LIGHT,
    "ytick.color":      HEX_LIGHT,
    "text.color":       HEX_WHITE,
    "grid.color":       HEX_GRID,
    "grid.linestyle":   "--",
    "grid.alpha":       0.5,
    "font.size":        11,
    "axes.titlesize":   13,
    "axes.titlecolor":  HEX_WHITE,
})


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150,
                facecolor=HEX_BG)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── pptx 헬퍼 ────────────────────────────────────────────────────────────────

def fill_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color=C_CARD, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def accent_bar(slide, top=Inches(0.12)):
    """슬라이드 상단 가로 액센트 바"""
    add_rect(slide, 0, top, SLIDE_W, Inches(0.06), C_ACCENT)


def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    add_text(slide, title,
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             size=28, bold=True, color=C_ACCENT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.9), Inches(12), Inches(0.45),
                 size=14, color=C_LIGHT)


def add_image_bytes(slide, img_bytes, left, top, width, height):
    slide.shapes.add_picture(io.BytesIO(img_bytes) if isinstance(img_bytes, bytes) else img_bytes,
                             left, top, width, height)


# ── 차트 생성 함수들 ──────────────────────────────────────────────────────────

def chart_honeypot_volume():
    data = {
        "suricata":    1544775,
        "p0f":         1703477,
        "cowrie":        92873,
        "honeytrap":     79090,
        "fatt":          78104,
        "dionaea":       23725,
        "sentrypeer":    18694,
        "tanner":         2569,
        "others":         4350,
    }
    labels = list(data.keys())
    values = list(data.values())

    colors = [HEX_ACC, HEX_ACC2, "#4CAF50", "#9C27B0", "#FF9800",
              "#E91E63", "#00BCD4", "#8BC34A", "#607D8B"]

    fig, ax = plt.subplots(figsize=(7, 4.5))
    wedges, texts, autotexts = ax.pie(
        values, labels=None, autopct="%1.1f%%",
        colors=colors, startangle=140,
        pctdistance=0.78,
        wedgeprops=dict(linewidth=1.5, edgecolor=HEX_BG)
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color(HEX_WHITE)
    ax.legend(wedges, [f"{l} ({v:,})" for l, v in zip(labels, values)],
              loc="center left", bbox_to_anchor=(1, 0.5),
              fontsize=8.5, frameon=False)
    ax.set_facecolor(HEX_BG)
    fig.patch.set_facecolor(HEX_BG)
    ax.set_title("허니팟별 이벤트 비중 (총 3,547,589건)", color=HEX_WHITE, pad=10)
    return fig_to_bytes(fig)


def chart_cowrie_top_ip():
    data = {
        "117.172.163.46":  12088,
        "104.156.52.218":   7238,
        "208.115.218.22":   5460,
        "104.243.42.83":    3694,
        "192.109.200.50":   2058,
        "212.67.214.156":   1614,
        "94.237.60.218":    1600,
        "64.235.40.106":    1480,
        "162.0.234.212":    1312,
        "206.189.150.94":   1242,
    }
    fig, ax = plt.subplots(figsize=(7, 4))
    ips = list(data.keys())[::-1]
    vals = [data[i] for i in ips]
    bars = ax.barh(ips, vals, color=HEX_ACC, height=0.6)
    for bar, v in zip(bars, vals):
        ax.text(v + 100, bar.get_y() + bar.get_height()/2,
                f"{v:,}", va="center", fontsize=9, color=HEX_WHITE)
    ax.set_xlabel("접속 시도 횟수", color=HEX_LIGHT)
    ax.set_title("SSH 브루트포스 Top 10 공격자 IP (cowrie)", color=HEX_WHITE)
    ax.set_xlim(0, max(vals) * 1.18)
    ax.grid(axis="x")
    return fig_to_bytes(fig)


def chart_cowrie_credentials():
    users = {"root": 2587, "admin": 37, "ubuntu": 23, "deploy": 15,
             "user": 10, "pi": 9, "test": 7, "guest": 6, "oracle": 5, "postgres": 4}
    pwds  = {"123456": 39, "admin": 30, "123": 14, "root": 12,
             "12345": 12, "password": 12, "1234": 8, "92113504": 8,
             "test": 7, "1234567890": 6}

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    for ax, d, title, color in [
        (ax1, users, "Top 10 사용자명", HEX_ACC),
        (ax2, pwds,  "Top 10 비밀번호",  HEX_ACC2),
    ]:
        keys = list(d.keys())[::-1]
        vals = [d[k] for k in keys]
        bars = ax.barh(keys, vals, color=color, height=0.6)
        for bar, v in zip(bars, vals):
            ax.text(v + 5, bar.get_y() + bar.get_height()/2,
                    str(v), va="center", fontsize=9, color=HEX_WHITE)
        ax.set_title(title, color=HEX_WHITE)
        ax.set_xlabel("시도 횟수", color=HEX_LIGHT)
        ax.grid(axis="x")
    fig.tight_layout(pad=2)
    return fig_to_bytes(fig)


def chart_cowrie_hourly():
    hours = list(range(24))
    counts = [1631,5644,2235,1920,679,825,658,1829,1679,2106,
              10485,2821,12189,4542,3341,2780,2430,4667,5074,3038,
              6688,1214,7815,6571]
    fig, ax = plt.subplots(figsize=(10, 3.5))
    ax.fill_between(hours, counts, alpha=0.35, color=HEX_ACC)
    ax.plot(hours, counts, color=HEX_ACC, lw=2.5, marker="o", markersize=4)
    ax.set_xticks(hours)
    ax.set_xlabel("시간 (UTC)", color=HEX_LIGHT)
    ax.set_ylabel("이벤트 수", color=HEX_LIGHT)
    ax.set_title("SSH 공격 시간대 분포 (cowrie, UTC)", color=HEX_WHITE)
    ax.grid(axis="y")
    peak_h = hours[counts.index(max(counts))]
    ax.axvline(peak_h, color=HEX_ACC2, linestyle="--", lw=1.5,
               label=f"피크 {peak_h}시 ({max(counts):,}건)")
    ax.legend(frameon=False)
    return fig_to_bytes(fig)


def chart_suricata_category():
    cats = {
        "Misc Attack":                              69015,
        "Misc activity":                            24950,
        "Potentially Bad Traffic":                  20341,
        "Attempted Admin Privilege Gain":           18481,
        "Attempted Information Leak":               11069,
        "Vulnerable Web Application":               10338,
        "Generic Protocol Decode":                   6804,
        "Network Scan":                               428,
        "Not Suspicious":                             427,
        "Web Application Attack":                     351,
    }
    fig, ax = plt.subplots(figsize=(8, 4.5))
    labels = list(cats.keys())[::-1]
    vals   = [cats[l] for l in labels]
    colors = [HEX_ACC if v > 10000 else HEX_ACC2 if v > 5000 else HEX_LIGHT
              for v in vals]
    bars = ax.barh(labels, vals, color=colors, height=0.65)
    for bar, v in zip(bars, vals):
        ax.text(v + 300, bar.get_y() + bar.get_height()/2,
                f"{v:,}", va="center", fontsize=9, color=HEX_WHITE)
    ax.set_xlabel("알림 수", color=HEX_LIGHT)
    ax.set_title("Suricata IDS 알림 카테고리 Top 10 (162,602건)", color=HEX_WHITE)
    ax.set_xlim(0, max(vals) * 1.18)
    ax.grid(axis="x")
    return fig_to_bytes(fig)


def chart_suricata_ports():
    ports = {445: 18454, 5900: 6394, 80: 2891, 8080: 1845,
             53: 1784, 22: 1237, 23: 1118, 443: 1038,
             5060: 838, 161: 807}
    labels = [f":{p}" for p in ports.keys()][::-1]
    vals   = list(ports.values())[::-1]
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.barh(labels, vals, color=HEX_ACC2, height=0.6)
    ax.set_title("공격 목표 포트 Top 10 (Suricata)", color=HEX_WHITE)
    ax.set_xlabel("알림 수", color=HEX_LIGHT)
    ax.grid(axis="x")
    return fig_to_bytes(fig)


def chart_dionaea_protocols():
    protos = {"SMB": 22392, "MySQL": 248, "MongoDB": 244, "MSSQL": 224,
              "PPTP": 181, "HTTP": 151, "FTP": 150, "EPMapper": 107,
              "MQTT": 16}
    labels = list(protos.keys())
    vals   = list(protos.values())
    colors_pie = [HEX_ACC, HEX_ACC2, "#4CAF50", "#9C27B0", "#FF9800",
                  "#E91E63", "#00BCD4", "#8BC34A", "#607D8B"]
    fig, ax = plt.subplots(figsize=(6, 4))
    wedges, texts, autotexts = ax.pie(
        vals, labels=labels, autopct="%1.1f%%",
        colors=colors_pie, startangle=90,
        wedgeprops=dict(linewidth=1.5, edgecolor=HEX_BG),
        textprops={"color": HEX_WHITE, "fontsize": 9}
    )
    for at in autotexts:
        at.set_fontsize(9)
    ax.set_facecolor(HEX_BG)
    fig.patch.set_facecolor(HEX_BG)
    ax.set_title("Dionaea 타겟 서비스 분포 (23,725건)", color=HEX_WHITE)
    return fig_to_bytes(fig)


def chart_sentrypeer():
    methods = {"INVITE": 14972, "REGISTER": 3478, "OPTIONS": 238, "기타": 2}
    fig, ax = plt.subplots(figsize=(5, 3.5))
    colors_s = [HEX_ACC, HEX_ACC2, "#4CAF50", HEX_LIGHT]
    ax.bar(methods.keys(), methods.values(), color=colors_s, width=0.5)
    for k, v in zip(methods.keys(), methods.values()):
        ax.text(list(methods.keys()).index(k), v + 150, f"{v:,}",
                ha="center", fontsize=10, color=HEX_WHITE)
    ax.set_title("VoIP 공격 SIP 메서드 분포 (SentryPeer)", color=HEX_WHITE)
    ax.set_ylabel("이벤트 수", color=HEX_LIGHT)
    ax.grid(axis="y")
    return fig_to_bytes(fig)


# ── 슬라이드 빌더 ──────────────────────────────────────────────────────────────

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(slide)

    # 중앙 큰 제목 박스
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4),
             C_CARD)
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(0.12), Inches(4),
             C_ACCENT)  # 왼쪽 세로 바

    add_text(slide, "T-Pot 허니팟 공격 데이터셋 분석",
             Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
             size=36, bold=True, color=C_ACCENT)
    add_text(slide, "수집 기간: 2026-04-27 ~ 2026-05-06  |  총 이벤트: 3,547,589건  |  18개 허니팟",
             Inches(1.2), Inches(3.0), Inches(11), Inches(0.6),
             size=16, color=C_LIGHT)
    add_text(slide, "AWS EC2 · T-Pot 24.04.1 · Elasticsearch · ML Classifier",
             Inches(1.2), Inches(3.7), Inches(11), Inches(0.6),
             size=13, color=RGBColor(0x60, 0x90, 0xB0))

    add_text(slide, "캡스톤 33팀",
             Inches(1.2), Inches(5.0), Inches(5), Inches(0.5),
             size=15, bold=True, color=C_WHITE)


def build_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    accent_bar(slide)
    slide_title(slide, "데이터셋 개요",
                "10일간 AWS EC2 T-Pot 허니팟에서 수집된 원시 공격 로그")

    # KPI 카드 4개
    kpis = [
        ("3,547,589", "총 이벤트"),
        ("18개",       "허니팟 종류"),
        ("1,383+",     "고유 공격자 IP"),
        ("10일",       "수집 기간"),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.6)
    gap    = Inches(0.35)
    start_x = Inches(0.6)
    top    = Inches(1.5)
    for i, (val, label) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, top, card_w, card_h, C_CARD)
        add_rect(slide, x, top, card_w, Inches(0.06), C_ACCENT)
        add_text(slide, val,   x, top + Inches(0.15), card_w, Inches(0.8),
                 size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, top + Inches(0.85), card_w, Inches(0.5),
                 size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # 허니팟 목록 테이블
    table_data = [
        ["허니팟",      "용도",               "이벤트 수"],
        ["suricata",   "네트워크 IDS",         "1,544,775"],
        ["p0f",        "OS 핑거프린팅",         "1,703,477"],
        ["cowrie",     "SSH/Telnet 허니팟",    "92,873"],
        ["honeytrap",  "TCP 포트 스캔 캡처",    "79,090"],
        ["fatt",       "프로토콜 핑거프린팅",    "78,104"],
        ["dionaea",    "멀웨어 수집",            "23,725"],
        ["sentrypeer", "VoIP/SIP 공격",         "18,694"],
        ["기타 11종",   "—",                    "~6,000"],
    ]
    col_w = [Inches(2.2), Inches(4.2), Inches(2.2)]
    row_h = Inches(0.36)
    tx = Inches(0.6)
    ty = Inches(3.3)
    for r, row in enumerate(table_data):
        cx = tx
        for c, (cell, w) in enumerate(zip(row, col_w)):
            is_hdr = (r == 0)
            bg = C_ACCENT if is_hdr else (C_CARD if r % 2 == 0 else RGBColor(0x10, 0x22, 0x33))
            add_rect(slide, cx, ty + r * row_h, w, row_h, bg)
            add_text(slide, cell, cx + Inches(0.08),
                     ty + r * row_h + Inches(0.04), w - Inches(0.1), row_h - Inches(0.06),
                     size=11 if not is_hdr else 12,
                     bold=is_hdr,
                     color=C_BG if is_hdr else C_WHITE)
            cx += w

    img = chart_honeypot_volume()
    add_image_bytes(slide, img, Inches(8.8), Inches(1.5), Inches(4.2), Inches(4.8))


def build_ssh(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    slide_title(slide, "SSH 브루트포스 공격 분석 (Cowrie)",
                "92,873건 SSH/Telnet 이벤트 · 41,725 세션 · 55건 로그인 성공")

    img1 = chart_cowrie_top_ip()
    add_image_bytes(slide, img1, Inches(0.3), Inches(1.5), Inches(6.2), Inches(3.5))

    # 통계 카드
    stats = [
        ("41,725",  "총 세션"),
        ("55",      "로그인 성공"),
        ("1,383",   "고유 IP"),
        ("12,088",  "최다 시도 IP"),
    ]
    for i, (v, l) in enumerate(stats):
        x = Inches(6.8) + i * Inches(1.55)
        add_rect(slide, x, Inches(1.5), Inches(1.45), Inches(1.1), C_CARD)
        add_rect(slide, x, Inches(1.5), Inches(1.45), Inches(0.05), C_ACCENT2)
        add_text(slide, v, x, Inches(1.6), Inches(1.45), Inches(0.55),
                 size=20, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, l, x, Inches(2.1), Inches(1.45), Inches(0.4),
                 size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    img2 = chart_cowrie_credentials()
    add_image_bytes(slide, img2, Inches(0.3), Inches(5.0), Inches(12.5), Inches(2.3))


def build_ssh_time(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    slide_title(slide, "SSH 공격 시간대 패턴 분석",
                "UTC 기준 시간대별 이벤트 분포 — 업무 시간대 집중 경향")

    img = chart_cowrie_hourly()
    add_image_bytes(slide, img, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0))

    insights = [
        "피크 12시 UTC: 12,189건 (한국 기준 오후 9시)",
        "10시·22시 UTC 에 보조 피크 — 자동화 스캐너 활동",
        "새벽 4~6시 UTC 최저 — 일부 지역 야간 시간대",
        "24시간 지속 공격 → 완전 자동화된 봇넷 활동",
    ]
    for i, ins in enumerate(insights):
        add_text(slide, f"• {ins}",
                 Inches(0.6), Inches(5.65) + i * Inches(0.38),
                 Inches(12), Inches(0.35),
                 size=12, color=C_LIGHT)


def build_suricata(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    slide_title(slide, "네트워크 IDS 알림 분석 (Suricata)",
                "162,602건 알림 · Severity 1~3 · 주요 공격 카테고리")

    img1 = chart_suricata_category()
    add_image_bytes(slide, img1, Inches(0.3), Inches(1.5), Inches(7.5), Inches(5.5))

    img2 = chart_suricata_ports()
    add_image_bytes(slide, img2, Inches(7.9), Inches(1.5), Inches(5.2), Inches(3.8))

    # severity 요약
    sev_data = [("Severity 1 (High)", "19,139건", C_ACCENT2),
                ("Severity 2 (Med)",  "110,852건", C_ACCENT),
                ("Severity 3 (Low)",  "32,611건",  C_LIGHT)]
    for i, (label, val, col) in enumerate(sev_data):
        add_rect(slide, Inches(7.9) + i * Inches(1.8), Inches(5.5),
                 Inches(1.7), Inches(0.9), C_CARD)
        add_text(slide, val,   Inches(7.9) + i * Inches(1.8), Inches(5.55),
                 Inches(1.7), Inches(0.45), size=16, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        add_text(slide, label, Inches(7.9) + i * Inches(1.8), Inches(5.95),
                 Inches(1.7), Inches(0.4), size=9, color=C_LIGHT,
                 align=PP_ALIGN.CENTER)


def build_service_attacks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    slide_title(slide, "서비스별 공격 분석 (Dionaea · SentryPeer · Heralding)",
                "멀웨어 수집, VoIP 공격, 인증 크리덴셜 수집")

    img1 = chart_dionaea_protocols()
    add_image_bytes(slide, img1, Inches(0.3), Inches(1.5), Inches(5.5), Inches(4.0))

    img2 = chart_sentrypeer()
    add_image_bytes(slide, img2, Inches(5.9), Inches(1.5), Inches(4.5), Inches(3.5))

    # heralding 프로토콜
    add_rect(slide, Inches(10.6), Inches(1.5), Inches(2.5), Inches(3.5), C_CARD)
    add_text(slide, "Heralding\n인증 시도 (127건)",
             Inches(10.6), Inches(1.55), Inches(2.5), Inches(0.7),
             size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    protocols_h = [("VNC",        91, HEX_ACC),
                   ("SOCKS5",     30, HEX_ACC2),
                   ("PostgreSQL",  6, "#4CAF50")]
    for i, (proto, cnt, col) in enumerate(protocols_h):
        y = Inches(2.35) + i * Inches(0.7)
        add_rect(slide, Inches(10.7), y, Inches(2.2), Inches(0.55), C_CARD)
        pct = cnt / 127
        bar_w = Inches(2.2) * pct
        add_rect(slide, Inches(10.7), y, max(bar_w, Inches(0.05)), Inches(0.55),
                 RGBColor(*bytes.fromhex(col[1:])))
        add_text(slide, f"{proto}  {cnt}건",
                 Inches(10.7), y + Inches(0.12), Inches(2.2), Inches(0.35),
                 size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 공통 인사이트
    insights = [
        "SMB(445) 집중 공격 → EternalBlue / 랜섬웨어 전파 시도",
        "VoIP INVITE 폭탄 → 무단 국제전화 도용 목적",
        "VNC/SOCKS5 크리덴셜 수집 → 프록시 체인 구성 의심",
    ]
    for i, ins in enumerate(insights):
        add_text(slide, f"• {ins}",
                 Inches(0.5), Inches(5.7) + i * Inches(0.38),
                 Inches(12.5), Inches(0.35),
                 size=12, color=C_LIGHT)


def build_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    slide_title(slide, "주요 발견 및 결론")

    findings = [
        ("1. 지속적 자동화 공격",
         "24시간 SSH 브루트포스, 1,383+ IP, 피크 UTC 12시 — 글로벌 봇넷 활동"),
        ("2. SMB(445) 집중 타겟",
         "Suricata 알림 중 445포트 18,454건 1위 — WannaCry 계열 랜섬웨어 스캔"),
        ("3. 단순 크리덴셜 재사용",
         "top 비밀번호: 123456·admin·password — 기본값/약한 비밀번호 위협"),
        ("4. VoIP 공격 다수",
         "SentryPeer 18,694건 INVITE 폭탄 — SIP 스캐너 봇 다수 활동"),
        ("5. SMB 멀웨어 수집 집중",
         "Dionaea SMB 연결 22,392건 (94%) — 파일 공유 취약점 자동 악용"),
    ]

    for i, (title, desc) in enumerate(findings):
        y = Inches(1.4) + i * Inches(1.05)
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.95), C_CARD)
        add_rect(slide, Inches(0.4), y, Inches(0.1), Inches(0.95),
                 C_ACCENT if i % 2 == 0 else C_ACCENT2)
        add_text(slide, title,
                 Inches(0.65), y + Inches(0.05), Inches(4.5), Inches(0.4),
                 size=13, bold=True,
                 color=C_ACCENT if i % 2 == 0 else C_ACCENT2)
        add_text(slide, desc,
                 Inches(0.65), y + Inches(0.45), Inches(11.8), Inches(0.45),
                 size=12, color=C_LIGHT)


# ── 실행 ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("슬라이드 생성 중...")
    print("  [1/6] 표지")
    build_cover(prs)
    print("  [2/6] 데이터셋 개요")
    build_overview(prs)
    print("  [3/6] SSH 브루트포스 분석")
    build_ssh(prs)
    print("  [4/6] 시간대 패턴")
    build_ssh_time(prs)
    print("  [5/6] Suricata IDS")
    build_suricata(prs)
    print("  [6/7] 서비스별 공격")
    build_service_attacks(prs)
    print("  [7/7] 주요 발견 및 결론")
    build_summary(prs)

    prs.save(PPT_PATH)
    size_mb = PPT_PATH.stat().st_size / 1024 / 1024
    print(f"\n✓  저장 완료: {PPT_PATH}  ({size_mb:.1f} MB)")


if __name__ == "__main__":
    main()
